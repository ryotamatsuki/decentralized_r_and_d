from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()
    
    # Titles and Slide content definitions
    slides_data = [
        {
            "title": "Decentralized Public R&D, Demand Shifters, and the Over- and Underinvestment Threshold\n(分権的公共R&D、需要シフトと過剰・過少投資の閾値)",
            "content": [
                "発表者: 松木 遼太 (Ryota Matsuki)",
                "応用経済学会"
            ],
            "is_title": True
        },
        {
            "title": "背景 (Background)",
            "content": [
                "地方自治体による特産品のブランド化・品質向上への投資",
                "コーヒー豆、農産物、観光地などのプロモーションや研究開発",
                "これらの投資は消費者の「知覚品質（Perceived Quality）」を向上・シフトさせる",
                "地域間の激しい競争：各自治体が自地域の製品を売るために競い合っている"
            ],
            "is_title": False
        },
        {
            "title": "問題意識 (Problem Statement)",
            "content": [
                "自治体の公共支出は「真の品質向上」をもたらしているのか？",
                "それとも、単なる他地域からの「需要の奪い合い（Business Stealing / Persuasion）」に過ぎないのか？",
                "説得・プロモーション競争に偏る場合、社会的無駄（Promotion-policy trap）になる恐れ",
                "一方で、知識スピルオーバーが存在する真のR&Dであれば、過少投資になる可能性"
            ],
            "is_title": False
        },
        {
            "title": "アプローチと研究の目的 (Approach & Objective)",
            "content": [
                "知覚品質（Perceived Quality）を2つの要素に分解:",
                "1. 実質的品質（Real Quality）: 厚生に寄与する割合 α",
                "2. 説得的要素（Persuasive Component）: 厚生には寄与せず需要のみをシフトさせる割合 1-α",
                "ホテリング・モデルを用いて自治体間のR&D競争を分析",
                "分権的R&Dが社会的に「過剰」になるか「過少」になるかの明確な閾値（Threshold）を導出"
            ],
            "is_title": False
        },
        {
            "title": "主要な貢献 (Main Contributions)",
            "content": [
                "1. 自治体間の分権的R&D投資における過剰・過少投資の解析的な閾値を導出",
                "2. 中央政府（または制度的調整メカニズム）による最適なマッチング補助率 (s*) の導出",
                "3. 真のR&Dとプロモーションを分ける2投入モデルへの拡張により、分権化が「プロモーション偏向（Promotion tilt）」を引き起こし、内生的に実質的品質の割合(α)を低下させることを証明"
            ],
            "is_title": False
        },
        {
            "title": "モデル設定1：空間競争 (Model Setup: Spatial Competition)",
            "content": [
                "2つの地域（自治体 i=1, 2）、およびそれぞれに対応する下流企業 A, B",
                "消費者は区間 [0,1]（ホテリング直線）上に一様分布",
                "各企業は限界費用 c で生産し、消費者は輸送費用 t に直面する",
                "消費者の効用:",
                "U = V + q_bar + q_i - p_i - t*d(x)",
                "ここで q_i は投資によって引き上げられる「知覚品質」"
            ],
            "is_title": False
        },
        {
            "title": "モデル設定2：品質投資とスピルオーバー (R&D & Spillovers)",
            "content": [
                "自治体 i が R&D投資 e_i を決定し、費用 K(e_i) を負担",
                "知識スピルオーバー（パラメータ ρ ∈ [0,1)）:",
                "企業Aの品質: q_A = g(e_1) + ρ*g(e_2)",
                "企業Bの品質: q_B = g(e_2) + ρ*g(e_1)",
                "自地域の投資だけでなく、他地域の投資の恩恵を部分的に受ける"
            ],
            "is_title": False
        },
        {
            "title": "社会的厚生基準 (Social Welfare Criterion)",
            "content": [
                "本研究における中核的な厚生評価の前提:",
                "知覚品質 q_f のうち、厚生に実質的に寄与するのは r_f = α*q_f のみとみなす",
                "残りの部分 (1-α)*q_f は、「純粋に説得的な要素（Purely persuasive component）」であり、厚生評価から除外",
                "最適化問題（社会計画者）は、この真の品質向上と、輸送・生産・投資費用のトレードオフを検討"
            ],
            "is_title": False
        },
        {
            "title": "ゲームのタイミング (Timing of the Game)",
            "content": [
                "第0段階: 中央政府がマッチング補助率 s を設定",
                "第1段階: 各自治体が R&D投資 e_1, e_2 を同時に決定",
                "第2段階: 下流企業 A, B が価格 p_A, p_B を同時に決定（ベルトラン競争）",
                "第3段階: 消費者が購入先を選択し、市場がクリアされる"
            ],
            "is_title": False
        },
        {
            "title": "第2段階：価格競争の均衡 (Stage 2: Price Equilibrium)",
            "content": [
                "品質差 Δq = q_A - q_B が 3t より小さい時、内部解が存在",
                "均衡価格:",
                "p_A* = c + t + Δq/3",
                "p_B* = c + t - Δq/3",
                "縮約形利潤 (Reduced-form profit):",
                "π_A*(q_A, q_B) = (3t + Δq)^2 / (18t)"
            ],
            "is_title": False
        },
        {
            "title": "第1段階：分権的R&D均衡 (Stage 1: Decentralized R&D Equilibrium)",
            "content": [
                "自治体の目的関数（マッチング率 s の下で）:",
                "W_i = π_i*(Δq) - (1-s)K(e_i)",
                "対称なナッシュ均衡 (e_1 = e_2 = e^N) での一階条件:",
                "(1-s)K'(e^N) = (1-ρ)/3 * g'(e^N)",
                "自治体は「自己の投資が自国利潤を増やす（Business stealing）」インセンティブに基づいて投資水準を決定"
            ],
            "is_title": False
        },
        {
            "title": "最適R&D水準 (Social Optimum)",
            "content": [
                "社会計画者は、実質的品質 (α) と全体のコストを最小化するように e^S を設定",
                "対称な社会的最適での一階条件:",
                "K'(e^S) = [α(1+ρ)/2] * g'(e^S)",
                "（スピルオーバー ρ がプラスに働く）"
            ],
            "is_title": False
        },
        {
            "title": "分析結果1：過剰・過少についての閾値定理",
            "content": [
                "分権的均衡(e^N) は、以下の条件に従って社会的最適(e^S)から乖離する:",
                "過剰投資になる条件: α < 2(1-ρ) / {3(1+ρ)}",
                "過少投資になる条件: α > 2(1-ρ) / {3(1+ρ)}",
                "閾値となる限界水準を α_bar(ρ) = 2(1-ρ) / {3(1+ρ)} と定義"
            ],
            "is_title": False
        },
        {
            "title": "閾値の直感とメカニズム (Intuition)",
            "content": [
                "過剰投資（α が小さい領域）:",
                "投資が主に「説得的・プロモーション的」であり、厚生に寄与しないが、自治体は需要を奪うために過剰な宣伝競争に陥る",
                "過少投資（α と ρ が大きい領域）:",
                "投資が有意義な「実質的品質向上」をもたらし、他地域へのスピルオーバーも大きいが、自治体は他国への恩恵を無視するため過少となる"
            ],
            "is_title": False
        },
        {
            "title": "分析結果2：マッチング補助金による実装",
            "content": [
                "中央政府は、単一の政策手段（補助率 s*）を用いて社会的最適を実装可能",
                "s* = 1 - α_bar(ρ) / α",
                "α < α_bar の場合（過剰投資領域）:",
                "s* は負の値となり、R&D支出に対する「税金」または「規制（上限設定）」が必要",
                "α > α_bar の場合（過少投資領域）:",
                "s* は正の値となり、適切な「補助金」が必要"
            ],
            "is_title": False
        },
        {
            "title": "拡張モデル：2投入プロセスとプロモーション偏向",
            "content": [
                "投資を「真のR&D（r_i）」と「プロモーション（m_i）」の2つに分解した拡張分析",
                "限界費用を k_R, k_M とし、スピルオーバーを ρ_R, ρ_M と仮定",
                "結果:",
                "自治体はプロモーション(m_i)に社会最適(=0)よりも過剰に投資し（Promotion-policy trap）、真のR&Dには過少投資する",
                "結果として、投資全体における真の品質の比率（内生的な α_N）が低下してしまう（Promotion tilt）"
            ],
            "is_title": False
        },
        {
            "title": "政策的含意 (Policy Implications)",
            "content": [
                "地方創生や農水産物の輸出促進において、自治体のプロモーション寄りの支出は「需要の奪い合い」になりやすく、社会的な無駄となる可能性がある",
                "一律の補助金支給ではなく、政策手段の性質（品質改善 vs 純粋なブランディング）を見極めることが重要",
                "実質的品質改善（品種改良や安全性保証システム）にターゲットを絞ったインセンティブ設計や、スピルオーバーの共有が不可欠"
            ],
            "is_title": False
        },
        {
            "title": "結論 (Conclusion)",
            "content": [
                "地域間のR&Dおよびブランディング競争において、「知覚品質の構成要素」を分解することで、過剰投資と過少投資の境界を解析的に明らかにした",
                "品質投資が「真の改善」か「需要の奪い合い」かによって、補助すべきか課税（規制）すべきかが明確に逆転する",
                "制度設計においては、プロモーション偏向を防ぐ仕組みが求められる"
            ],
            "is_title": False
        }
    ]

    for item in slides_data:
        if item["is_title"]:
            slide_layout = prs.slide_layouts[0] # Title Slide
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = item["title"]
            subtitle.text = "\\n".join(item["content"])
            # Adjust font sizes
            for p in title.text_frame.paragraphs:
                p.font.size = Pt(28)
            for p in subtitle.text_frame.paragraphs:
                p.font.size = Pt(20)
        else:
            slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = item["title"]
            text_frame = content.text_frame
            for i, line in enumerate(item["content"]):
                if i == 0:
                    text_frame.text = line
                else:
                    p = text_frame.add_paragraph()
                    p.text = line
            
            for p in text_frame.paragraphs:
                p.font.size = Pt(20)
                
    prs.save('presentation_applied_econ.pptx')
    print("presentation_applied_econ.pptx created successfully.")

if __name__ == '__main__':
    create_presentation()
